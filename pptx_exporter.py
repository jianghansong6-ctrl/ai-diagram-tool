import json
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE

# Map canvas (1200x800) to slide (10"x7.5"), no extra offsets
SCALE = 10 / 1200  # 1 canvas unit = ~0.00833 inches


def to_inch(val: float, default: float = 100) -> float:
    """Convert canvas coordinate to inches."""
    return (val if val is not None else default) * SCALE


def _get_zIndex(inst: dict) -> int:
    """Extract zIndex from an instruction, handling string or dict params."""
    params = inst.get("params", {})
    if isinstance(params, str):
        try:
            params = json.loads(params)
        except Exception:
            return 0
    return params.get("zIndex", 0)


def instructions_to_pptx(instructions: list[dict], output_path: str, title: str = "Scientific Diagram"):
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    slide_layout = prs.slide_layouts[6]  # blank layout
    slide = prs.slides.add_slide(slide_layout)


    # Sort by zIndex to match canvas rendering order
    sorted_instructions = sorted(instructions, key=_get_zIndex)

    # Pre-pass: collect all shape bounding boxes for overlap detection
    all_shape_bounds: list[tuple] = []
    for inst in instructions:
        action = inst.get("action", "")
        params = inst.get("params", {})
        if isinstance(params, str):
            params = json.loads(params)
        bounds = _get_shape_bounds(params, action)
        if bounds:
            all_shape_bounds.append(bounds)

    for inst in sorted_instructions:
        action = inst.get("action", "")
        params = inst.get("params", {})
        if isinstance(params, str):
            params = json.loads(params)

        try:
            shape = None
            x = to_inch(params.get("x"))
            y = to_inch(params.get("y"))

            if action == "draw_rect":
                w = to_inch(params.get("w", 100))
                h = to_inch(params.get("h", 100))
                # Use rounded rectangle if rx > 0
                rx = params.get("rx", 0)
                shape_id = MSO_SHAPE.ROUNDED_RECTANGLE if rx else MSO_SHAPE.RECTANGLE
                shape = slide.shapes.add_shape(shape_id, Inches(x), Inches(y), Inches(w), Inches(h))
                _apply_fill(shape, params.get("fill", "#CCCCCC"), params.get("stroke", "#333333"), params.get("strokeWidth", 2))

            elif action == "draw_circle":
                r = to_inch(params.get("r", 30))
                shape = slide.shapes.add_shape(
                    MSO_SHAPE.OVAL,
                    Inches(x - r), Inches(y - r),
                    Inches(r * 2), Inches(r * 2)
                )
                _apply_fill(shape, params.get("fill", "#CCCCCC"), params.get("stroke", "#333333"), params.get("strokeWidth", 2))

            elif action == "draw_ellipse":
                # Canvas uses (x,y) as center; PPTX uses top-left
                ex = to_inch(params.get("x"))
                ey = to_inch(params.get("y"))
                ew = to_inch(params.get("w", 80))
                eh = to_inch(params.get("h", 40))
                shape = slide.shapes.add_shape(
                    MSO_SHAPE.OVAL,
                    Inches(ex - ew / 2), Inches(ey - eh / 2),
                    Inches(ew), Inches(eh)
                )
                _apply_fill(shape, params.get("fill", "#CCCCCC"), params.get("stroke", "#333333"), params.get("strokeWidth", 2))

            elif action == "draw_line" or action == "draw_dashed_line":
                sx = to_inch(params.get("startX", 0))
                sy = to_inch(params.get("startY", 0))
                ex = to_inch(params.get("endX", 100))
                ey = to_inch(params.get("endY", 100))
                connector = slide.shapes.add_connector(1, Inches(sx), Inches(sy), Inches(ex), Inches(ey))
                connector.line.color.rgb = RGBColor(*hex_to_rgb(params.get("stroke", "#666666")))
                connector.line.width = Pt(params.get("strokeWidth", 1))

            elif action == "draw_arrow":
                sx = to_inch(params.get("startX", 0))
                sy = to_inch(params.get("startY", 0))
                ex = to_inch(params.get("endX", 100))
                ey = to_inch(params.get("endY", 100))
                midX = params.get("midX")
                midY = params.get("midY")
                if midX is not None:
                    mx = to_inch(midX)
                    _draw_line_segment(slide, sx, sy, mx, sy, params.get("stroke", "#333333"), 2, False)
                    _draw_line_segment(slide, mx, sy, mx, ey, params.get("stroke", "#333333"), 2, False)
                    _draw_line_segment(slide, mx, ey, ex, ey, params.get("stroke", "#333333"), 2, True)
                elif midY is not None:
                    my = to_inch(midY)
                    _draw_line_segment(slide, sx, sy, sx, my, params.get("stroke", "#333333"), 2, False)
                    _draw_line_segment(slide, sx, my, ex, my, params.get("stroke", "#333333"), 2, False)
                    _draw_line_segment(slide, ex, my, ex, ey, params.get("stroke", "#333333"), 2, True)
                else:
                    _draw_line_segment(slide, sx, sy, ex, ey, params.get("stroke", "#333333"), 2, True)

            elif action == "draw_pie_slice":
                cx = to_inch(params.get("x", 0)) + to_inch(params.get("r", 50))
                cy = to_inch(params.get("y", 0)) + to_inch(params.get("r", 50))
                size = to_inch(params.get("r", 50) * 2)
                shape = slide.shapes.add_shape(
                    MSO_SHAPE.CHORD if shape_id else MSO_SHAPE.OVAL,
                    Inches(cx - size / 2), Inches(cy - size / 2),
                    Inches(size), Inches(size)
                )
                shape.fill.solid()
                shape.fill.fore_color.rgb = RGBColor(*hex_to_rgb(params.get("fill", "#3B82F6")))
                # Label inside pie slice
                if params.get("label"):
                    label_tf = shape.text_frame
                    label_tf.word_wrap = True
                    label_tf.paragraphs[0].text = params["label"]
                    label_tf.paragraphs[0].font.size = Pt(12)
                    label_tf.paragraphs[0].alignment = 2
                    label_tf.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)

            elif action == "draw_text":
                text_str = params.get("text", "")
                # Estimate available width: distance from x to right edge (~10")
                avail_w = max(2.0, (10 - x - 0.3))
                orig_size = params.get("fontSize", 14)
                fitted_size = _fit_font_size(text_str, avail_w, 2.0, orig_size)
                txBox = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(avail_w), Inches(0.5))
                tf = txBox.text_frame
                tf.word_wrap = True
                p = tf.paragraphs[0]
                p.text = text_str
                p.font.size = Pt(fitted_size)
                p.font.color.rgb = RGBColor(*hex_to_rgb(params.get("fontColor", "#2C3E50")))
                ta = params.get("textAlign", "left")
                if ta == "center":
                    p.alignment = 2
                elif ta == "right":
                    p.alignment = 3
                else:
                    p.alignment = 1

            # Handle label on shapes
            label = params.get("label", "")
            if label and shape:
                # Use shape's actual EMU dimensions to fit font
                sw = shape.width / 914400  # EMU → inches
                sh = shape.height / 914400
                orig_size = params.get("fontSize", 13)
                fitted_size = _fit_font_size(label, sw, sh, orig_size)
                tf = shape.text_frame
                tf.word_wrap = True
                tf.paragraphs[0].text = label
                tf.paragraphs[0].font.size = Pt(fitted_size)
                tf.paragraphs[0].alignment = 2
                # If another shape overlaps this one, push text to bottom
                this_bounds = _get_shape_bounds(params, action)
                if this_bounds:
                    has_overlap = any(
                        _rects_overlap(this_bounds, b)
                        for b in all_shape_bounds
                        if b != this_bounds
                    )
                    if has_overlap:
                        # Set vertical anchor to bottom so text stays visible
                        tf._txBody.bodyPr.set('anchor', 'b')
                # Auto white/dark text based on fill luminance
                label_color = params.get("fontColor", "")
                if not label_color:
                    fill = params.get("fill", "#CCCCCC")
                    r, g, b = hex_to_rgb(fill)
                    lum = r * 0.299 + g * 0.587 + b * 0.114
                    label_color = "#FFFFFF" if lum < 128 else "#1A1A2E"
                tf.paragraphs[0].font.color.rgb = RGBColor(*hex_to_rgb(label_color))
                tf.paragraphs[0].space_before = Pt(0)
                tf.paragraphs[0].space_after = Pt(0)

        except Exception:
            continue

    prs.save(output_path)
    return output_path


def _apply_fill(shape, fill_color: str, stroke_color: str, stroke_width: int):
    """Apply fill and stroke to a shape."""
    fill = shape.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(*hex_to_rgb(fill_color))
    shape.line.color.rgb = RGBColor(*hex_to_rgb(stroke_color))
    shape.line.width = Pt(stroke_width)


def hex_to_rgb(hex_color: str) -> tuple:
    hex_color = hex_color.lstrip("#")
    if len(hex_color) == 3:
        hex_color = "".join(c * 2 for c in hex_color)
    return tuple(int(hex_color[i:i+2], 16) for i in (0, 2, 4))


def _fit_font_size(text: str, w_inches: float, h_inches: float, original_size: int, min_size: int = 6) -> int:
    """Reduce font size so text fits within box bounds in PPTX (fonts differ from Canvas)."""
    if not text or w_inches <= 0 or h_inches <= 0:
        return original_size
    # Estimate: avg char width ≈ 0.65 × fontSize (pts); 1pt = 1/72 inch.
    # Conservative for mixed CJK + Latin text.
    est_char_width = 0.65
    max_by_w = int(w_inches * 72 / (len(text) * est_char_width)) if len(text) > 0 else original_size
    # Height: single line needs fontSize + ~4pt overhead
    max_by_h = int((h_inches * 72 - 4) / 1.3)
    return max(min_size, min(original_size, max_by_w, max_by_h))


def _get_shape_bounds(params: dict, action: str) -> tuple[float, float, float, float] | None:
    """Get (left, top, right, bottom) in inches for a shape-creating action, or None."""
    x = to_inch(params.get("x"))
    y = to_inch(params.get("y"))
    if action == "draw_rect":
        w = to_inch(params.get("w", 100))
        h = to_inch(params.get("h", 100))
        return (x, y, x + w, y + h)
    elif action == "draw_circle":
        r = to_inch(params.get("r", 30))
        return (x - r, y - r, x + r, y + r)
    elif action == "draw_ellipse":
        ew = to_inch(params.get("w", 80))
        eh = to_inch(params.get("h", 40))
        return (x - ew / 2, y - eh / 2, x + ew / 2, y + eh / 2)
    elif action == "draw_pie_slice":
        r = to_inch(params.get("r", 50))
        return (x, y, x + r * 2, y + r * 2)
    return None


def _rects_overlap(a: tuple, b: tuple) -> bool:
    """Check if two bounding rectangles (left, top, right, bottom) overlap."""
    return not (a[2] <= b[0] or a[0] >= b[2] or a[3] <= b[1] or a[1] >= b[3])


def _draw_line_segment(slide, x1, y1, x2, y2, color, width, arrowhead=False):
    """Draw a single connector line segment in PPTX."""
    connector = slide.shapes.add_connector(1, Inches(x1), Inches(y1), Inches(x2), Inches(y2))
    connector.line.color.rgb = RGBColor(*hex_to_rgb(color))
    connector.line.width = Pt(width)
    if arrowhead:
        ln = connector.line._element.find('{http://schemas.openxmlformats.org/drawingml/2006/main}ln')
        if ln is not None:
            ln.attrib['tailEnd'] = 'triangle'
